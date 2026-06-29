"""PPTXRenderer:DeckSpec → .pptx 字节。确定性、零 LLM。

这是唯一出现 EMU 的地方:归一化 Rect × 画布 EMU(12192000×6858000)。Design Token
(font_role→字号、color_role→调色板、fonts)在此落地。渲染由 slot_mapper 产出的
RenderSlot 驱动,renderer 只按 kind 分派,不关心 slide_type 细节。

素材(AssetSpec)可选:有 background binding 且文件存在则贴全幅底图 + 可选半透明遮罩,
并用 safe_area 把该页文字收进安全区;文件缺失则安全跳过(回退纯色)。
"""

from __future__ import annotations

import re
from collections.abc import Mapping
from contextlib import suppress
from dataclasses import replace
from io import BytesIO
from math import ceil
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

from packages.asset_engine import AssetEngine
from packages.core import (
    AssetBinding,
    AssetSpec,
    CardsSlide,
    DataSlide,
    DeckSpec,
    Rect,
    SlideSpec,
    ThemeSpec,
)
from packages.core.enums import ColorRole, FontRole
from packages.layout_engine import LayoutEngine, map_slots
from packages.layout_engine.slot_mapper import IconResolver, RenderSlot, SlotValue

# 图标徽章按角色取色:卡片图标压在 primary 圆上用 on_primary(=背景色)、流程图标用 accent。
_ICON_COLOR_ROLE = {"card": "on_primary", "process": "accent"}

# 16:9 画布的 EMU 尺寸(13.333in × 7.5in)
CANVAS_W = 12192000
CANVAS_H = 6858000
_DEFAULT_ASSETS_DIR = Path(__file__).resolve().parents[2] / "assets"

_ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
_VALIGN = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}
_HEADING_ROLES = ("display", "h1", "h2")
_MIN_HEADING_PT = 18.0  # display/h1 标题自适应缩字号下限,低于此宁可截断也不再缩
_MIN_SUBHEAD_PT = 14.0  # h2(卡片标题/对比小标/单列小结条)自适应下限
_MIN_BODY_PT = 11.0  # 正文/要点(卡片正文/对比要点/流程·时间轴说明)自适应下限
# 这些 deck_type 的封面用居中变体(克制/正式风)
_COVER_CENTERED_TYPES = frozenset({"exec_report", "pitch", "tech"})


def _cover_variant(slide_type: str, deck_type: str) -> str | None:
    """按 deck_type 决定封面布局变体:正式类型用居中封面,其余用默认。"""
    if slide_type == "cover" and deck_type in _COVER_CENTERED_TYPES:
        return "cover_centered"
    return None


def _rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color.lstrip("#"))


def _confine(rect: Rect, safe: Rect) -> Rect:
    """把 [0,1] 布局矩形线性映射进 safe_area 子框,保证文字落在背景图安全区内。"""
    return Rect(
        left=safe.left + rect.left * safe.width,
        top=safe.top + rect.top * safe.height,
        width=rect.width * safe.width,
        height=rect.height * safe.height,
    )


def _truncate(text: str, limit: int) -> str:
    """超出 limit 字数则截断 + 省略号,避免文字溢出框。"""
    return text if len(text) <= limit else text[: max(1, limit - 1)] + "…"


def _text_em(text: str) -> float:
    """文本的"字宽当量"(em):CJK/全角≈1.0、ASCII≈0.55、空格≈0.3。比固定系数更贴近真实换行。"""
    total = 0.0
    for ch in text:
        o = ord(ch)
        if ch == " ":
            total += 0.3
        elif o < 0x2000:  # ASCII / Latin
            total += 0.55
        elif o < 0x3000:  # 通用标点
            total += 0.6
        else:  # CJK 及全角
            total += 1.0
    return total


_TEXTBOX_MARGIN = 182880  # 文本框默认左右内边距合计(2×0.1in)
_WRAP_SAFETY = 1.35  # 换行保守系数:宁可偏小,确保实际渲染不超估行宽


def _per_line_em(rect: Rect, size: float) -> float:
    """该字号下,rect 一行能容纳的"字宽当量"(em)数。"""
    line_w = max(1.0, rect.width * CANVAS_W - _TEXTBOX_MARGIN)
    return max(1.0, line_w / (size * 12700))


def _max_rows(rect: Rect, size: float, line_height: float) -> int:
    """该字号下 rect 纵向能容纳的行数(行高=字号×line_height)。"""
    return max(1, int(rect.height * CANVAS_H / (size * 12700 * line_height)))


def _wrap_lines(text: str, per_line_em: float, bullet_em: float = 0.0) -> int:
    """文本按字宽当量估算换行行数(含可选项目符宽度,保守取整向上)。"""
    em = (_text_em(text) + bullet_em) * _WRAP_SAFETY
    return max(1, ceil(em / per_line_em))


def _fit_text_size(
    value: SlotValue,
    rect: Rect,
    base: float,
    line_height: float,
    min_pt: float,
    *,
    bullet_em: float = 0.0,
    step: float = 1.0,
) -> float:
    """正文/标题按字宽当量 + 框宽高自适应字号:从 base 逐级缩小,直到全部内容在框内
    不溢出(总换行数 ≤ 可容行数),最低到 min_pt。str 当单段;list 每项独立换行。

    bullet_em 计入条目首行的项目符宽度(列表用)。这是 _fit_heading_size 的推广,标题、
    卡片正文、对比要点、流程/时间轴说明共用同一套估算,长文先缩字号、缩到底再由 _fit_lines 截断。
    """
    items = [it for it in (value if isinstance(value, list) else [value]) if it]
    if not items:
        return base
    size = base
    while size > min_pt:
        per_line = _per_line_em(rect, size)
        total = sum(_wrap_lines(it, per_line, bullet_em) for it in items)
        if total <= _max_rows(rect, size, line_height):
            return size
        size -= step
    return min_pt


def _fit_heading_size(text: str, rect: Rect, base: float, line_height: float) -> float:
    """标题自适应字号(_fit_text_size 的标题特化:下限 18pt、步进 2pt)。"""
    return _fit_text_size(text, rect, base, line_height, _MIN_HEADING_PT, step=2.0)


def _fit_lines(
    value: SlotValue, rect: Rect, size: float, line_height: float, bullet: str
) -> list[str]:
    """给定字号,产出要画的段落(每个 list 项一段,可换行);按框高 max_rows 截断,
    超出末项按剩余行数估字数截断 + 省略号。自适应已把字号缩到尽量放下,这里只兜底防溢出。"""
    per_line = _per_line_em(rect, size)
    rows = _max_rows(rect, size, line_height)
    bullet_em = _text_em(bullet)
    items = value if isinstance(value, list) else [value]
    out: list[str] = []
    used = 0
    for item in items:
        if used >= rows:
            break
        need = _wrap_lines(item, per_line, bullet_em)
        if used + need <= rows:
            out.append(bullet + item)
            used += need
        else:
            remain = rows - used
            avg_em = max(0.3, _text_em(item) / max(1, len(item)))  # 每字 em 当量
            budget = remain * per_line / _WRAP_SAFETY - bullet_em
            char_cap = max(1, int(budget / avg_em))
            out.append(bullet + _truncate(item, char_cap))
            used = rows
    return out


_NUM_RE = re.compile(r"-?\d+(?:\.\d+)?")
# 后缀→倍率(按"先长后短/避免误命中"排序:亿/万 在前)
_SCALES = (("亿", 1e8), ("万", 1e4), ("K", 1e3), ("k", 1e3), ("M", 1e6), ("m", 1e6), ("B", 1e9), ("b", 1e9))


def _parse_number(text: str) -> float | None:
    """取首个数值并按后缀(K/M/B/万/亿)换算;'78%'→78、'¥1.2M'→1.2e6、'1.2万'→12000;无数字则 None。"""
    cleaned = text.replace(",", "")
    m = _NUM_RE.search(cleaned)
    if m is None:
        return None
    val = float(m.group())
    tail = cleaned[m.end() :]
    for suffix, mul in _SCALES:
        if suffix in tail:
            return val * mul
    return val


def _data_chart_values(slide: SlideSpec) -> list[tuple[str, float]] | None:
    """DataSlide 且每个 metric.value 可解析、≥2 个 → (标签,数值) 列表;否则 None(回退文本块)。"""
    if not isinstance(slide, DataSlide):
        return None
    pairs: list[tuple[str, float]] = []
    for m in slide.metrics:
        num = _parse_number(m.value)
        if num is None:
            return None
        pairs.append((m.label or m.value, num))
    return pairs if len(pairs) >= 2 else None


class PPTXRenderer:
    def __init__(
        self, layout_engine: LayoutEngine | None = None, assets_dir: Path | None = None
    ) -> None:
        self._layouts = layout_engine or LayoutEngine()
        self._assets_dir = Path(assets_dir) if assets_dir else _DEFAULT_ASSETS_DIR
        # 图标解析:读预制 Lucide 图标库(零栅格化);未构建则 resolve 恒返回 None → 回退徽章。
        self._assets = AssetEngine(self._assets_dir / "asset_index.json")

    def render(
        self,
        deck: DeckSpec,
        assets: AssetSpec | None = None,
        *,
        layout_map: Mapping[str, str] | None = None,
    ) -> bytes:
        """渲染整套 deck。layout_map(可选,来自选定模板的 preset.layout)按页型覆盖版式变体。"""
        prs = Presentation()
        prs.slide_width = Emu(CANVAS_W)
        prs.slide_height = Emu(CANVAS_H)
        blank = prs.slide_layouts[6]  # 空白版式
        backgrounds = {
            b.slide_id: b for b in (assets.bindings if assets else []) if b.role == "background"
        }
        total = len(deck.slides)
        for i, slide in enumerate(deck.slides):
            self._render_slide(
                prs,
                blank,
                slide,
                deck.theme,
                backgrounds.get(slide.id),
                page=i + 1,
                total=total,
                deck_title=deck.title,
                deck_type=deck.deck_type,
                layout_map=layout_map,
            )
        buf = BytesIO()
        prs.save(buf)
        return buf.getvalue()

    # ── 内部 ──────────────────────────────────────────────────────
    def _render_slide(
        self,
        prs: Any,
        blank: Any,
        slide: SlideSpec,
        theme: ThemeSpec,
        binding: AssetBinding | None,
        *,
        page: int,
        total: int,
        deck_title: str,
        deck_type: str,
        layout_map: Mapping[str, str] | None = None,
    ) -> None:
        s = prs.slides.add_slide(blank)
        # 书挡/章节页用浅色 tint 底(surface,浅而不平),内容页纯背景;有底图时不覆盖
        bookend = slide.type in ("cover", "section", "big_idea", "closing")
        bg = theme.palette.surface if (bookend and binding is None) else theme.palette.background
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = _rgb(bg)

        safe: Rect | None = None
        if binding is not None:
            image = self._assets_dir / binding.object_key
            if image.exists():
                s.shapes.add_picture(str(image), Emu(0), Emu(0), Emu(CANVAS_W), Emu(CANVAS_H))
                if binding.needs_overlay and binding.overlay_opacity > 0:
                    self._draw_scrim(s, binding.overlay_opacity)
                safe = binding.safe_area

        # 母题:封面/章节/大主张/结尾左侧贯穿色脊(primary),跨页一致;不在标题下/旁加装饰横线
        if slide.type in ("cover", "section", "big_idea", "closing"):
            self._draw_motif(s, theme, Rect(left=0.0, top=0.0, width=0.014, height=1.0), "primary")
        # 封面/结尾:右下大号柔光圆 + 右上语义图标,填留白、立母题(柔光圆在文字下,不挡读)
        if slide.type in ("cover", "closing"):
            self._soft_circle(s, theme.palette.primary, self._square_rect(0.92, 0.9, 0.66), 8)
            kw = getattr(slide, "title", "") or getattr(slide, "subtitle", "") or getattr(slide, "cta", "")
            self._draw_motif_icon(
                s, theme, kw, self._square_rect(0.85, 0.2, 0.12), "primary", fallback="sparkles"
            )
        # 章节页:右侧大号语义图标 + 柔光圆(始终在场,语义 miss 退默认,不再右半空)
        if slide.type == "section":
            self._soft_circle(s, theme.palette.primary, self._square_rect(0.76, 0.5, 0.46), 8)
            self._draw_motif_icon(
                s, theme, slide.title, self._square_rect(0.76, 0.5, 0.32), "primary", fallback="sparkles"
            )
        # big_idea 字少:右侧大号半透明 accent 母题块 + 左上语义图标,填充视觉留白(不挡居中文字)
        if slide.type == "big_idea":
            self._draw_accent_field(s, theme)
            statement = getattr(slide, "statement", "")
            self._draw_motif_icon(
                s, theme, statement, self._square_rect(0.135, 0.245, 0.09), "accent", fallback="sparkles"
            )

        count = len(slide.cards) if isinstance(slide, CardsSlide) else 0
        # 选版式:模板 layout 映射优先(换模板换版式变体);模板没指定该页型 → deck_type 兜底
        # (如正式封面居中)。变体名→真实 layout JSON 的解析由 LayoutEngine 持有。
        layout = self._layouts.pick(
            slide.type,
            content_count=count,
            layout_overrides=layout_map,
            default_variant=_cover_variant(slide.type, deck_type),
        )
        chart_pairs = _data_chart_values(slide)
        if layout.id == "data_big_number":
            chart_pairs = None  # big_number 变体:强制大号数字 stat callout,不画柱状图
        chart_region: Rect | None = None
        for rs in map_slots(slide, layout, self._icon_for(theme)):
            placed = replace(rs, rect=_confine(rs.rect, safe)) if safe is not None else rs
            if chart_pairs is not None and rs.name.startswith("metric_"):
                chart_region = self._union(chart_region, placed.rect)  # 合并指标区→图表区
                continue
            if slide.type == "cover" and rs.name == "kicker" and placed.align != "center":
                placed = self._draw_kicker(s, theme, placed)  # 左对齐封面:accent 小方块 + 右移文字
            if placed.kind == "image":
                self._draw_icon(s, placed)
            elif placed.kind == "shape":
                self._draw_surface(s, placed, theme)
            else:
                self._draw_text(s, placed, theme)
        if chart_pairs is not None and chart_region is not None:
            self._draw_chart(s, chart_region, chart_pairs, theme)

        # 页脚/页码:cover/closing 不加;branded 额外带 deck 标题
        if theme.footer_style != "none" and slide.type not in ("cover", "closing"):
            title = deck_title if theme.footer_style == "branded" else ""
            self._draw_footer(s, theme, page, total, title)

        # 演讲者备注:写进 pptx 备注页(演示时可见,投屏不显)
        if slide.speaker_notes:
            s.notes_slide.notes_text_frame.text = slide.speaker_notes

    def _draw_chart(
        self, s: Any, rect: Rect, pairs: list[tuple[str, float]], theme: ThemeSpec
    ) -> None:
        """在 rect 内画简洁原生柱状图(无图例/标题,柱色用 accent,带数据标签)。"""
        data = CategoryChartData()
        data.categories = [p[0] for p in pairs]
        data.add_series("", [p[1] for p in pairs])
        gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, *self._emu(rect), data)
        chart = gf.chart
        chart.has_legend = False
        chart.has_title = False
        text_rgb = _rgb(theme.palette.text)
        # 隐藏数值轴 + 网格线:更干净,且避免大数值左轴标签/数据标签换行;
        # 类目轴与柱顶数据标签用主题文字色(暗色主题也可见,不再黑字压暗底)。
        val_ax = chart.value_axis
        val_ax.visible = False
        val_ax.has_major_gridlines = False
        cat_ax = chart.category_axis
        cat_ax.has_major_gridlines = False
        cat_ax.tick_labels.font.size = Pt(12)
        cat_ax.tick_labels.font.name = theme.fonts.body
        cat_ax.tick_labels.font.color.rgb = text_rgb
        with suppress(Exception):
            cat_ax.format.line.color.rgb = _rgb(theme.palette.border)
        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.gap_width = 70
        dl = plot.data_labels
        dl.font.size = Pt(12)
        dl.font.bold = True
        dl.font.name = theme.fonts.body
        dl.font.color.rgb = text_rgb
        for series in plot.series:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = _rgb(theme.palette.accent)

    @staticmethod
    def _union(a: Rect | None, b: Rect) -> Rect:
        """两个归一化 Rect 的包围盒(把若干 metric 槽合成一块图表区域)。"""
        if a is None:
            return b
        left, top = min(a.left, b.left), min(a.top, b.top)
        right = max(a.left + a.width, b.left + b.width)
        bottom = max(a.top + a.height, b.top + b.height)
        return Rect(left=left, top=top, width=right - left, height=bottom - top)

    def _draw_motif(self, s: Any, theme: ThemeSpec, rect: Rect, role: str) -> None:
        """实心母题块(色脊/小方块):无描边无阴影,色取自 palette。"""
        shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, *self._emu(rect))
        shp.fill.solid()
        shp.fill.fore_color.rgb = self._color(role, theme)  # type: ignore[arg-type]
        shp.line.fill.background()
        shp.shadow.inherit = False

    def _draw_accent_field(self, s: Any, theme: ThemeSpec) -> None:
        """big_idea 这类字少页的视觉母题:右侧一枚大号半透明 accent 圆,部分出血,填留白不挡文字。"""
        rect = Rect(left=0.74, top=0.12, width=0.46, height=0.46 * (CANVAS_W / CANVAS_H))
        shp = s.shapes.add_shape(MSO_SHAPE.OVAL, *self._emu(rect))
        shp.fill.solid()
        shp.fill.fore_color.rgb = _rgb(theme.palette.accent)
        shp.line.fill.background()
        shp.shadow.inherit = False
        with suppress(Exception):
            self._set_alpha(shp, 12)

    def _soft_circle(self, s: Any, hex_color: str, rect: Rect, alpha_pct: int) -> None:
        """半透明柔光圆(母题装饰,压在文字下不挡读);alpha 写不进则安全跳过。"""
        shp = s.shapes.add_shape(MSO_SHAPE.OVAL, *self._emu(rect))
        shp.fill.solid()
        shp.fill.fore_color.rgb = _rgb(hex_color)
        shp.line.fill.background()
        shp.shadow.inherit = False
        with suppress(Exception):
            self._set_alpha(shp, alpha_pct)

    @staticmethod
    def _square_rect(cx: float, cy: float, size_h: float) -> Rect:
        """以 (cx,cy) 为中心、视觉上正方形的归一化 Rect(按画布宽高比收窄宽度)。"""
        w = size_h * (CANVAS_H / CANVAS_W)
        return Rect(left=cx - w / 2, top=cy - size_h / 2, width=w, height=size_h)

    def _draw_motif_icon(
        self,
        s: Any,
        theme: ThemeSpec,
        keyword: str,
        rect: Rect,
        role: str,
        fallback: str | None = None,
    ) -> None:
        """按关键词派生一枚语义图标贴进 rect(role 取色,预制 PNG)。

        语义无命中且给了 fallback(图标名)→ 直接取该图标的对应色版,保证母题始终在场。
        无图标库 / 颜色未预制 / 文件缺失 → 安全跳过(页面仍有色脊等母题)。
        """
        color = self._color(role, theme)  # type: ignore[arg-type]
        key = self._assets.resolve_icon(keyword, f"{color}") if keyword else None
        image = self._assets_dir / key if key else None
        if (image is None or not image.exists()) and fallback:
            cand = self._assets_dir / "icons" / "lucide" / f"{fallback}__{str(color).lower()}.png"
            image = cand if cand.exists() else None
        if image is not None and image.exists():
            s.shapes.add_picture(str(image), *self._emu(rect))

    def _draw_kicker(self, s: Any, theme: ThemeSpec, rs: RenderSlot) -> RenderSlot:
        """封面 kicker:左侧画一个 accent 小方块(母题点),文字右移让位。"""
        r = rs.rect
        dot = Rect(left=r.left, top=r.top + r.height * 0.18, width=0.012, height=r.height * 0.5)
        self._draw_motif(s, theme, dot, "accent")
        shifted = Rect(left=r.left + 0.026, top=r.top, width=r.width - 0.026, height=r.height)
        return replace(rs, rect=shifted)

    def _draw_footer(self, s: Any, theme: ThemeSpec, page: int, total: int, title: str) -> None:
        self._footer_text(
            s,
            theme,
            Rect(left=0.78, top=0.93, width=0.16, height=0.05),
            f"{page} / {total}",
            "right",
        )
        if title:
            self._footer_text(
                s, theme, Rect(left=0.06, top=0.93, width=0.55, height=0.05), title, "left"
            )

    def _footer_text(self, s: Any, theme: ThemeSpec, rect: Rect, text: str, align: str) -> None:
        box = s.shapes.add_textbox(*self._emu(rect))
        para = box.text_frame.paragraphs[0]
        para.alignment = _ALIGN[align]
        run = para.add_run()
        run.text = _truncate(text, 60)
        run.font.size = Pt(theme.type_scale.caption)
        run.font.name = theme.fonts.body
        run.font.color.rgb = _rgb(theme.palette.text_muted)

    def _draw_scrim(self, s: Any, opacity: float) -> None:
        """在底图上压一层黑色半透明遮罩,提升文字对比度。alpha 写不进则安全跳过。"""
        try:
            shp = s.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Emu(CANVAS_W), Emu(CANVAS_H)
            )
            shp.line.fill.background()
            shp.shadow.inherit = False
            shp.fill.solid()
            shp.fill.fore_color.rgb = RGBColor(0, 0, 0)
            self._set_alpha(shp, int(opacity * 100))
        except Exception:  # noqa: BLE001 — 遮罩是锦上添花,失败不应阻断渲染
            pass

    @staticmethod
    def _set_alpha(shape: Any, opacity_pct: int) -> None:
        from pptx.oxml.ns import qn

        pct = max(0, min(100, opacity_pct))
        solid = shape.fill._xPr.find(qn("a:solidFill"))
        srgb = solid.find(qn("a:srgbClr"))
        srgb.append(srgb.makeelement(qn("a:alpha"), {"val": str(pct * 1000)}))

    def _draw_surface(self, s: Any, rs: RenderSlot, theme: ThemeSpec) -> None:
        if rs.shape_role == "accent_bar":
            # 卡片左侧母题强调条:primary 实心圆角细条,无描边无阴影
            shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, *self._emu(rs.rect))
            self._round(shp, 0.5)
            shp.fill.solid()
            shp.fill.fore_color.rgb = _rgb(theme.palette.primary)
            shp.line.fill.background()
            shp.shadow.inherit = False
            return
        if rs.shape_role == "card_badge":
            # 卡片圆形徽章(母题):primary 实心圆,编号/图标以背景色压在其上
            shp = s.shapes.add_shape(MSO_SHAPE.OVAL, *self._emu(rs.rect))
            shp.fill.solid()
            shp.fill.fore_color.rgb = _rgb(theme.palette.primary)
            shp.line.fill.background()
            shp.shadow.inherit = False
            return
        if rs.shape_role == "num_badge":
            # 流程编号徽章:浅底圆 + accent 细描边,accent 数字落其内,对比清晰
            shp = s.shapes.add_shape(MSO_SHAPE.OVAL, *self._emu(rs.rect))
            shp.fill.solid()
            shp.fill.fore_color.rgb = _rgb(theme.palette.surface)
            shp.line.color.rgb = _rgb(theme.palette.accent)
            shp.line.width = Pt(1.25)
            shp.shadow.inherit = False
            return
        if rs.shape_role == "axis_line":
            # 时间轴主轴:细圆角条,色取 color_role(primary),压在节点下
            shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, *self._emu(rs.rect))
            self._round(shp, 0.5)
            shp.fill.solid()
            shp.fill.fore_color.rgb = self._color(rs.color_role, theme)
            shp.line.fill.background()
            shp.shadow.inherit = False
            return
        if rs.shape_role == "rule":
            # 流程连接线:border 色细线(克制,不与编号徽章争视觉)
            shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, *self._emu(rs.rect))
            shp.fill.solid()
            shp.fill.fore_color.rgb = _rgb(theme.palette.border)
            shp.line.fill.background()
            shp.shadow.inherit = False
            return
        if rs.shape_role == "node_dot":
            # 时间轴节点:实心圆(色取 color_role)+ 背景色细环(halo),落在轴线上
            shp = s.shapes.add_shape(MSO_SHAPE.OVAL, *self._emu(rs.rect))
            shp.fill.solid()
            shp.fill.fore_color.rgb = self._color(rs.color_role, theme)
            shp.line.color.rgb = _rgb(theme.palette.background)
            shp.line.width = Pt(2.5)
            shp.shadow.inherit = False
            return
        # card / metric / plain:浅彩底卡面 + 细边 + 轻阴影,与背景拉开
        shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, *self._emu(rs.rect))
        self._round(shp, 0.06)
        shp.fill.solid()
        shp.fill.fore_color.rgb = _rgb(theme.palette.surface)
        shp.line.color.rgb = _rgb(theme.palette.border)
        shp.line.width = Pt(0.75)
        shp.shadow.inherit = False
        self._soft_shadow(shp, theme)

    @staticmethod
    def _round(shape: Any, frac: float) -> None:
        """设置圆角矩形的圆角半径(相对短边的比例),避免默认过圆。"""
        with suppress(IndexError, KeyError):
            shape.adjustments[0] = frac

    @staticmethod
    def _soft_shadow(shape: Any, theme: ThemeSpec) -> None:
        """给形状加一层克制的外阴影(深色文字色 + 极低不透明度),提升卡片层次。"""
        from pptx.oxml.ns import qn

        spPr = shape._element.spPr
        for old in spPr.findall(qn("a:effectLst")):
            spPr.remove(old)
        eff = spPr.makeelement(qn("a:effectLst"), {})
        shdw = eff.makeelement(
            qn("a:outerShdw"),
            {"blurRad": "90000", "dist": "30000", "dir": "5400000", "rotWithShape": "0"},
        )
        clr = shdw.makeelement(qn("a:srgbClr"), {"val": theme.palette.text.lstrip("#")})
        clr.append(clr.makeelement(qn("a:alpha"), {"val": "13000"}))
        shdw.append(clr)
        eff.append(shdw)
        spPr.append(eff)

    def _icon_for(self, theme: ThemeSpec) -> IconResolver:
        """构造图标解析回调:按角色取主题色,委托 AssetEngine 解析成预制图标 PNG 的 object_key。"""

        def resolve(keyword: str, role: str) -> str | None:
            role_name = _ICON_COLOR_ROLE.get(role, "on_primary")
            color = self._color(role_name, theme)  # type: ignore[arg-type]
            return self._assets.resolve_icon(keyword, f"{color}")

        return resolve

    def _draw_icon(self, s: Any, rs: RenderSlot) -> None:
        """把预制图标 PNG 贴进 rect(透明底,叠在徽章圆上);文件缺失则安全跳过(回退徽章)。"""
        image = self._assets_dir / str(rs.value)
        if image.exists():
            s.shapes.add_picture(str(image), *self._emu(rs.rect))

    def _draw_text(self, s: Any, rs: RenderSlot, theme: ThemeSpec) -> None:
        box = s.shapes.add_textbox(*self._emu(rs.rect))
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = _VALIGN.get(rs.valign, MSO_ANCHOR.TOP)
        line_height = theme.type_scale.line_height
        bullet = "• " if rs.kind == "list" else ""
        # 按框自适应字号:正文/要点逐级缩到放得下(下限 11pt),标题/小标用各自下限。
        size = self._fit_size(rs, self._font_size(rs.font_role, theme), line_height)
        color = self._color(rs.color_role, theme)
        font_name = theme.fonts.heading if rs.font_role in _HEADING_ROLES else theme.fonts.body
        bold = rs.font_role in ("display", "h1")
        for i, line in enumerate(_fit_lines(rs.value, rs.rect, size, line_height, bullet)):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.alignment = _ALIGN.get(rs.align, PP_ALIGN.LEFT)
            run = para.add_run()
            run.text = line
            run.font.size = Pt(size)
            run.font.name = font_name
            run.font.color.rgb = color
            run.font.bold = bold

    @staticmethod
    def _fit_size(rs: RenderSlot, base: float, line_height: float) -> float:
        """按字号角色选下限/步进,调 _fit_text_size 自适应:display/h1→18pt、h2→14pt、
        正文/caption→11pt;列表项计入项目符宽度。空值/非文本直接用 base。"""
        if not rs.value:
            return base
        bullet_em = _text_em("• ") if rs.kind == "list" else 0.0
        if rs.font_role in ("display", "h1"):
            if not isinstance(rs.value, str):
                return base
            return _fit_text_size(rs.value, rs.rect, base, line_height, _MIN_HEADING_PT, step=2.0)
        min_pt = _MIN_SUBHEAD_PT if rs.font_role == "h2" else _MIN_BODY_PT
        return _fit_text_size(
            rs.value, rs.rect, base, line_height, min_pt, bullet_em=bullet_em
        )

    @staticmethod
    def _emu(rect: Any) -> tuple[Any, Any, Any, Any]:
        return (
            Emu(int(rect.left * CANVAS_W)),
            Emu(int(rect.top * CANVAS_H)),
            Emu(int(rect.width * CANVAS_W)),
            Emu(int(rect.height * CANVAS_H)),
        )

    @staticmethod
    def _font_size(role: FontRole | None, theme: ThemeSpec) -> float:
        ts = theme.type_scale
        sizes = {
            "display": ts.display,
            "h1": ts.h1,
            "h2": ts.h2,
            "body": ts.body,
            "caption": ts.caption,
        }
        return sizes.get(role or "body", ts.body)

    @staticmethod
    def _color(role: ColorRole, theme: ThemeSpec) -> RGBColor:
        p = theme.palette
        mapping = {
            "text": p.text,
            "text_muted": p.text_muted,
            "primary": p.primary,
            "accent": p.accent,
            "on_primary": p.background,
        }
        return _rgb(mapping.get(role, p.text))


def render_deck(deck: DeckSpec, assets: AssetSpec | None = None) -> bytes:
    """便捷入口:用默认 LayoutEngine 渲染一个 deck(可选素材)。"""
    return PPTXRenderer().render(deck, assets)
