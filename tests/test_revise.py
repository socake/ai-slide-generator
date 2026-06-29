from __future__ import annotations

from io import BytesIO

import pytest
from pptx import Presentation

from packages.pipeline import generate, render_pptx, revise_slide
from packages.planner import GenerationInput


def _deck() -> object:
    res = generate(
        GenerationInput(topic="Python 入门", brief="变量、循环、函数、列表", audience="零基础"),
        with_preview=False,
    )
    return res.deck_spec


def _first_cards_index(deck: object) -> int:
    return next(i for i, s in enumerate(deck.slides) if s.type == "cards")  # type: ignore[attr-defined]


def test_revise_changes_target_keeps_theme_and_others() -> None:
    deck = _deck()
    idx = _first_cards_index(deck)
    revised = revise_slide(deck, idx)  # type: ignore[arg-type]
    assert len(revised.slides) == len(deck.slides)  # type: ignore[attr-defined]
    assert revised.theme == deck.theme  # type: ignore[attr-defined]  全局主题不变
    assert revised.slides[idx].type != "cards"  # 轮换换型
    other = idx + 1
    assert revised.slides[other].title == deck.slides[other].title  # type: ignore[attr-defined]
    assert revised.slides[other].type == deck.slides[other].type  # type: ignore[attr-defined]


def test_revise_instruction_to_comparison() -> None:
    deck = _deck()
    idx = _first_cards_index(deck)
    revised = revise_slide(deck, idx, "换成对比页")  # type: ignore[arg-type]
    assert revised.slides[idx].type == "comparison"


def test_revise_simplify_keeps_type() -> None:
    deck = _deck()
    idx = _first_cards_index(deck)
    revised = revise_slide(deck, idx, "更简洁")  # type: ignore[arg-type]
    assert revised.slides[idx].type == "cards"
    assert len(revised.slides[idx].cards) <= 4


def test_revise_out_of_range() -> None:
    deck = _deck()
    with pytest.raises(IndexError):
        revise_slide(deck, 999)  # type: ignore[arg-type]


def test_revised_deck_renders() -> None:
    deck = _deck()
    revised = revise_slide(deck, 5, "加数据")  # type: ignore[arg-type]
    prs = Presentation(BytesIO(render_pptx(revised)))
    assert 25 <= len(prs.slides) <= 30
