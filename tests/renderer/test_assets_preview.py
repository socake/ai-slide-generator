from __future__ import annotations

from io import BytesIO
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from packages.core import (
    AssetBinding,
    AssetSpec,
    CoverSlide,
    DeckSpec,
    Fonts,
    NarrativeSpec,
    Palette,
    Rect,
    ThemeSpec,
)
from packages.renderer import PPTXRenderer, preview_available, to_png


def _deck_with_cover() -> DeckSpec:
    theme = ThemeSpec(
        id="t",
        name="t",
        mood=[],
        palette=Palette(
            primary="#111111",
            secondary="#222222",
            accent="#333333",
            background="#ffffff",
            surface="#f0f0f0",
            text="#000000",
            text_muted="#666666",
            border="#cccccc",
        ),
        fonts=Fonts(heading="A", body="B"),
    )
    return DeckSpec(
        id="d",
        title="T",
        topic="t",
        brief="b",
        audience="a",
        purpose="p",
        deck_type="generic",
        narrative=NarrativeSpec(hook="h", conflict="c", progression=["x"], resolution="r"),
        theme=theme,
        slides=[CoverSlide(id="s0", index=0, title="封面")],
    )


def test_background_image_is_embedded(tmp_path: Path) -> None:
    Image.new("RGB", (48, 27), (10, 20, 30)).save(tmp_path / "bg.png")
    assets = AssetSpec(
        bindings=[
            AssetBinding(
                slide_id="s0",
                role="background",
                object_key="bg.png",
                safe_area=Rect(left=0.1, top=0.5, width=0.8, height=0.4),
                needs_overlay=True,
                overlay_opacity=0.3,
            )
        ]
    )
    data = PPTXRenderer(assets_dir=tmp_path).render(_deck_with_cover(), assets)
    slide = Presentation(BytesIO(data)).slides[0]
    assert any(sh.shape_type == MSO_SHAPE_TYPE.PICTURE for sh in slide.shapes)


def test_missing_background_file_is_safe(tmp_path: Path) -> None:
    assets = AssetSpec(
        bindings=[AssetBinding(slide_id="s0", role="background", object_key="does_not_exist.png")]
    )
    data = PPTXRenderer(assets_dir=tmp_path).render(_deck_with_cover(), assets)
    # 文件不存在 → 不贴图、不抛;仍正常渲染
    assert len(Presentation(BytesIO(data)).slides) == 1


def test_preview_graceful_when_tools_absent() -> None:
    pngs = to_png(b"PK-not-a-real-pptx")
    assert pngs == [] or preview_available()  # 无 soffice/pdftoppm → 空列表,不抛
