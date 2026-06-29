from __future__ import annotations

from io import BytesIO

from pptx import Presentation
from pptx.enum.text import PP_ALIGN

from packages.core import (
    ClosingSlide,
    CoverSlide,
    DeckSpec,
    Fonts,
    NarrativeSpec,
    Palette,
    SlideSpec,
    ThemeSpec,
)
from packages.renderer import PPTXRenderer
from packages.renderer.renderer import _cover_variant


def _deck(deck_type: str) -> DeckSpec:
    theme = ThemeSpec(
        id="t",
        name="t",
        mood=[],
        palette=Palette(
            primary="#111111",
            secondary="#222222",
            accent="#333333",
            background="#ffffff",
            surface="#eeeeee",
            text="#000000",
            text_muted="#666666",
            border="#cccccc",
        ),
        fonts=Fonts(heading="A", body="B"),
        footer_style="none",
    )
    slides: list[SlideSpec] = [
        CoverSlide(id="c", index=0, title="封面", subtitle="副标题"),
        ClosingSlide(id="e", index=1, title=""),
    ]
    return DeckSpec(
        id="d",
        title="T",
        topic="t",
        brief="b",
        audience="a",
        purpose="p",
        deck_type=deck_type,
        narrative=NarrativeSpec(hook="h", conflict="c", progression=["x"], resolution="r"),
        theme=theme,
        slides=slides,
    )


def _center_count(prs_slide: object) -> int:
    n = 0
    for sh in prs_slide.shapes:  # type: ignore[attr-defined]
        if sh.has_text_frame:
            n += sum(
                1 for p in sh.text_frame.paragraphs if p.text and p.alignment == PP_ALIGN.CENTER
            )
    return n


def test_cover_variant_logic() -> None:
    assert _cover_variant("cover", "exec_report") == "cover_centered"
    assert _cover_variant("cover", "pitch") == "cover_centered"
    assert _cover_variant("cover", "teaching") is None
    assert _cover_variant("cards", "exec_report") is None


def test_formal_deck_cover_centered() -> None:
    cover = Presentation(BytesIO(PPTXRenderer().render(_deck("exec_report")))).slides[0]
    assert _center_count(cover) >= 1


def test_default_deck_cover_left() -> None:
    cover = Presentation(BytesIO(PPTXRenderer().render(_deck("teaching")))).slides[0]
    assert _center_count(cover) == 0
