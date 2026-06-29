from __future__ import annotations

from io import BytesIO

import pytest
from pptx import Presentation

from packages.core import DataSlide, DeckSpec, NarrativeSpec
from packages.core.slides import Metric
from packages.planner._defaults import neutral_theme
from packages.renderer import PPTXRenderer
from packages.renderer.renderer import _data_chart_values, _parse_number


@pytest.mark.parametrize(
    ("text", "expected"),
    [
        ("78%", 78.0),
        ("3.2x", 3.2),
        ("1.2万", 12000.0),
        ("¥1.2M", 1_200_000.0),
        ("800K", 800_000.0),
        ("2.5亿", 250_000_000.0),
        ("-5%", -5.0),
        ("1,200", 1200.0),
    ],
)
def test_parse_number_scales(text: str, expected: float) -> None:
    assert _parse_number(text) == expected


def test_parse_number_none_for_plain_text() -> None:
    assert _parse_number("高速增长") is None


def _data(metrics: list[Metric]) -> DataSlide:
    return DataSlide(id="d", index=0, title="指标", metrics=metrics)


def test_chart_values_parse_numeric() -> None:
    pairs = _data_chart_values(
        _data([Metric(value="78%", label="转化"), Metric(value="3.2x", label="增长")])
    )
    assert pairs == [("转化", 78.0), ("增长", 3.2)]


def test_chart_values_none_when_non_numeric() -> None:
    assert _data_chart_values(_data([Metric(value="高", label="a"), Metric(value="低", label="b")])) is None


def test_chart_values_none_when_single_metric() -> None:
    assert _data_chart_values(_data([Metric(value="78%", label="只一个")])) is None


def test_numeric_data_slides_render_native_charts() -> None:
    # 数值型 data 页出原生图表;非数值 data 页不出图(不编造数字、不强行画图)
    slides = [
        _data([Metric(value="¥1.2亿", label="营收"), Metric(value="38%", label="同比增长")]),
        _data([Metric(value="高", label="利润"), Metric(value="稳", label="留存")]),
    ]
    for i, s in enumerate(slides):
        s.index = i
        s.id = f"d{i}"
    deck = DeckSpec(
        id="deck-data", title="业绩", topic="t", brief="b", audience="a", purpose="p",
        deck_type="exec_report",
        narrative=NarrativeSpec(hook="h", conflict="c", progression=["x"], resolution="r"),
        theme=neutral_theme(),
        slides=slides,
    )
    expected = sum(1 for s in deck.slides if _data_chart_values(s) is not None)
    assert expected == 1  # 仅第一页可量化 → 仅它出图

    prs = Presentation(BytesIO(PPTXRenderer().render(deck)))
    charts = sum(
        1 for sl in prs.slides for sh in sl.shapes if getattr(sh, "has_chart", False)
    )
    assert charts == expected
