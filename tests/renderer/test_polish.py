from __future__ import annotations

from io import BytesIO

from pptx import Presentation

from packages.core import (
    BigIdeaSlide,
    ClosingSlide,
    CoverSlide,
    DeckSpec,
    Fonts,
    NarrativeSpec,
    Palette,
    SlideSpec,
    SummarySlide,
    ThemeSpec,
)
from packages.renderer import PPTXRenderer


def _theme(footer: str = "minimal") -> ThemeSpec:
    return ThemeSpec(
        id="t",
        name="t",
        mood=[],
        palette=Palette(
            primary="#2563eb",
            secondary="#1d4ed8",
            accent="#f59e0b",
            background="#ffffff",
            surface="#f1f5f9",
            text="#0f172a",
            text_muted="#64748b",
            border="#e2e8f0",
        ),
        fonts=Fonts(heading="A", body="B"),
        footer_style=footer,  # type: ignore[arg-type]
    )


def _deck(slides: list[SlideSpec], theme: ThemeSpec) -> DeckSpec:
    return DeckSpec(
        id="d",
        title="演示稿标题",
        topic="t",
        brief="b",
        audience="a",
        purpose="p",
        deck_type="generic",
        narrative=NarrativeSpec(hook="h", conflict="c", progression=["x"], resolution="r"),
        theme=theme,
        slides=slides,
    )


def _all_texts(prs_slide: object) -> list[str]:
    return [sh.text_frame.text for sh in prs_slide.shapes if sh.has_text_frame]  # type: ignore[attr-defined]


def _three() -> list[SlideSpec]:
    return [
        CoverSlide(id="c", index=0, title="封面"),
        SummarySlide(id="s", index=1, title="小结", points=["a"]),
        ClosingSlide(id="e", index=2, title=""),
    ]


def test_long_text_truncated_no_crash() -> None:
    huge = "很长的文本内容" * 300
    slides: list[SlideSpec] = [
        CoverSlide(id="c", index=0, title=huge, subtitle=huge),
        BigIdeaSlide(id="b", index=1, title=huge, statement=huge, support=huge),
        SummarySlide(id="s", index=2, title="小结", points=[huge, huge]),
        ClosingSlide(id="e", index=3, title=""),
    ]
    prs = Presentation(BytesIO(PPTXRenderer().render(_deck(slides, _theme()))))
    assert len(prs.slides) == 4
    for sl in prs.slides:
        for text in _all_texts(sl):
            assert len(text) < 800  # 被截断,不会整段塞进框


def test_footer_minimal_shows_page_number() -> None:
    prs = Presentation(BytesIO(PPTXRenderer().render(_deck(_three(), _theme("minimal")))))
    assert "2 / 3" in " ".join(_all_texts(prs.slides[1]))  # summary 第 2 页


def test_footer_none_hides_page_number() -> None:
    prs = Presentation(BytesIO(PPTXRenderer().render(_deck(_three(), _theme("none")))))
    assert "/" not in " ".join(t for sl in prs.slides for t in _all_texts(sl))


def test_footer_branded_shows_title() -> None:
    prs = Presentation(BytesIO(PPTXRenderer().render(_deck(_three(), _theme("branded")))))
    joined = " ".join(_all_texts(prs.slides[1]))
    assert "演示稿标题" in joined and "2 / 3" in joined
