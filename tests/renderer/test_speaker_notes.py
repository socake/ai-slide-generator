from __future__ import annotations

from io import BytesIO

from pptx import Presentation

from packages.pipeline import generate
from packages.planner import GenerationInput


def test_offline_content_slides_have_notes_and_render_to_pptx() -> None:
    res = generate(
        GenerationInput(topic="Python 入门", brief="变量、循环", audience="零基础"),
        with_preview=False,
    )
    deck = res.deck_spec
    expected = next((s.speaker_notes for s in deck.slides if s.speaker_notes), None)
    assert expected is not None  # 离线内容页应带 speaker_notes

    assert res.pptx_bytes is not None
    prs = Presentation(BytesIO(res.pptx_bytes))
    rendered = [sl.notes_slide.notes_text_frame.text for sl in prs.slides if sl.has_notes_slide]
    assert expected in rendered  # 备注已写进 pptx 备注页
