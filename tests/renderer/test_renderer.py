from __future__ import annotations

from io import BytesIO

from pptx import Presentation

from packages.core import (
    AgendaSlide,
    BigIdeaSlide,
    Card,
    CardsSlide,
    ClosingSlide,
    ComparisonColumn,
    ComparisonSlide,
    CoverSlide,
    DataSlide,
    DeckSpec,
    Metric,
    NarrativeSpec,
    ProcessSlide,
    ProcessStep,
    QuoteSlide,
    SectionSlide,
    SlideSpec,
    SummarySlide,
    TimelineEvent,
    TimelineSlide,
)
from packages.llm import MockLLMProvider
from packages.planner import GenerationInput, Planner
from packages.renderer import CANVAS_W, PPTXRenderer, render_deck
from packages.theme_engine import ThemeEngine


def _all_type_slides() -> list[SlideSpec]:
    return [
        CoverSlide(id="s0", index=0, title="封面", subtitle="副标题", kicker="2026"),
        AgendaSlide(id="s1", index=1, title="目录", items=["一", "二", "三"]),
        SectionSlide(id="s2", index=2, title="第一章", section_number=1),
        BigIdeaSlide(id="s3", index=3, title="主张", statement="一句话主张", support="支撑"),
        CardsSlide(
            id="s4",
            index=4,
            title="卡片",
            cards=[Card(title="A", body="a"), Card(title="B", body="b")],
        ),
        TimelineSlide(
            id="s5",
            index=5,
            title="时间轴",
            events=[TimelineEvent(time=f"Q{i}", title=f"E{i}") for i in range(1, 4)],
        ),
        ComparisonSlide(
            id="s6",
            index=6,
            title="对比",
            left=ComparisonColumn(heading="旧", points=["x"]),
            right=ComparisonColumn(heading="新", points=["y", "z"]),
        ),
        ProcessSlide(
            id="s7", index=7, title="流程", steps=[ProcessStep(title=f"S{i}") for i in range(3)]
        ),
        DataSlide(
            id="s8",
            index=8,
            title="数据",
            metrics=[
                Metric(value="78%", label="覆盖", delta="+5%"),
                Metric(value="3x", label="提速"),
            ],
        ),
        QuoteSlide(id="s9", index=9, title="", quote="一句引用", attribution="某人"),
        SummarySlide(id="s10", index=10, title="小结", points=["a", "b"]),
        ClosingSlide(id="s11", index=11, title="", subtitle="结束", cta="谢谢"),
    ]


def _deck(slides: list[SlideSpec], deck_type: str = "teaching") -> DeckSpec:
    return DeckSpec(
        id="d",
        title="T",
        topic="t",
        brief="b",
        audience="a",
        purpose="p",
        deck_type=deck_type,
        narrative=NarrativeSpec(hook="h", conflict="c", progression=["x"], resolution="r"),
        theme=ThemeEngine().select(deck_type),
        slides=slides,
    )


def test_render_every_slide_type_roundtrips() -> None:
    deck = _deck(_all_type_slides())
    data = PPTXRenderer().render(deck)
    assert data[:2] == b"PK"  # pptx 是 zip
    prs = Presentation(BytesIO(data))
    assert len(prs.slides) == 12
    assert prs.slide_width == CANVAS_W  # 16:9 画布


def test_render_planner_offline_deck() -> None:
    deck = Planner(MockLLMProvider()).generate(
        GenerationInput(topic="Python 入门", brief="变量/循环/函数", audience="零基础")
    )
    data = render_deck(deck)
    prs = Presentation(BytesIO(data))
    assert len(prs.slides) == len(deck.slides)
    assert 25 <= len(prs.slides) <= 30


def test_dark_theme_renders() -> None:
    # exec_report 是深色主题,确保深底也能渲染不抛
    deck = _deck(_all_type_slides(), deck_type="exec_report")
    prs = Presentation(BytesIO(PPTXRenderer().render(deck)))
    assert len(prs.slides) == 12
