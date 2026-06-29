from __future__ import annotations

from io import BytesIO

from pptx import Presentation

from packages.core import (
    AgendaSlide,
    BigIdeaSlide,
    Card,
    CardsSlide,
    ClosingSlide,
    ComparisonColumn,
    ComparisonSlide,
    CoverSlide,
    DataSlide,
    DeckSpec,
    Metric,
    NarrativeSpec,
    ProcessSlide,
    ProcessStep,
    QuoteSlide,
    SectionSlide,
    SlideSpec,
    SummarySlide,
    TimelineEvent,
    TimelineSlide,
)
from packages.planner._defaults import neutral_theme
from packages.renderer import PPTXRenderer


def _one_of_each() -> list[SlideSpec]:
    return [
        CoverSlide(id="s0", index=0, title="封面", subtitle="副标题", kicker="2026"),
        AgendaSlide(id="s1", index=1, title="目录", items=["一", "二", "三"]),
        SectionSlide(id="s2", index=2, title="第一章", section_number=1),
        BigIdeaSlide(id="s3", index=3, title="主张", statement="少即是多", support="支撑句"),
        CardsSlide(
            id="s4", index=4, title="要点",
            cards=[Card(title="A", body="aa"), Card(title="B", body="bb")],
        ),
        TimelineSlide(
            id="s5", index=5, title="时间线",
            events=[TimelineEvent(time="Q1", title="启动"), TimelineEvent(time="Q2", title="推进")],
        ),
        ComparisonSlide(
            id="s6", index=6, title="对比",
            left=ComparisonColumn(heading="左", points=["l1"]),
            right=ComparisonColumn(heading="右", points=["r1"]),
        ),
        ProcessSlide(
            id="s7", index=7, title="流程",
            steps=[ProcessStep(title="调研"), ProcessStep(title="设计")],
        ),
        DataSlide(
            id="s8", index=8, title="数据",
            metrics=[Metric(value="78%", label="转化"), Metric(value="3.2x", label="增长")],
        ),
        QuoteSlide(id="s9", index=9, title="金句", quote="大道至简", attribution="老子"),
        SummarySlide(id="s10", index=10, title="总结", points=["p1", "p2"]),
        ClosingSlide(id="s11", index=11, title="谢谢", subtitle="Q&A", cta="联系我们"),
    ]


def test_every_slide_type_renders() -> None:
    slides = _one_of_each()
    deck = DeckSpec(
        id="deck-all", title="全类型", topic="t", brief="b", audience="a", purpose="p",
        deck_type="generic",
        narrative=NarrativeSpec(hook="h", conflict="c", progression=["x"], resolution="r"),
        theme=neutral_theme(),
        slides=slides,
    )
    prs = Presentation(BytesIO(PPTXRenderer().render(deck)))
    assert len(prs.slides) == 12
    charts = sum(1 for sl in prs.slides for sh in sl.shapes if getattr(sh, "has_chart", False))
    assert charts == 1  # 唯一的数值 data 页出图
