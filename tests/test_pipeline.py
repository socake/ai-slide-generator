from __future__ import annotations

from io import BytesIO

from pptx import Presentation

from packages.pipeline import generate
from packages.planner import GenerationInput


def test_pipeline_generates_openable_deck() -> None:
    res = generate(GenerationInput(topic="Python 入门", brief="变量/循环/函数", audience="零基础"))
    assert res.pptx_bytes is not None
    assert res.pptx_bytes[:2] == b"PK"
    assert 25 <= len(res.deck_spec.slides) <= 30
    assert res.benchmark.slide_count == len(res.deck_spec.slides)
    assert res.benchmark.deck_id == res.deck_spec.id
    prs = Presentation(BytesIO(res.pptx_bytes))
    assert len(prs.slides) == len(res.deck_spec.slides)


def test_pipeline_applies_theme_engine() -> None:
    res = generate(GenerationInput(topic="季度业绩复盘", brief="b", audience="面向高管董事会"))
    assert res.deck_spec.theme.id == "exec_report"  # 高管受众 → 深色 exec 主题


def test_pipeline_is_deterministic() -> None:
    inp = GenerationInput(topic="同一主题", brief="同一简介", audience="同一受众")
    a = generate(inp)
    b = generate(inp)
    assert a.deck_spec.id == b.deck_spec.id
    assert a.deck_spec.model_dump() == b.deck_spec.model_dump()
