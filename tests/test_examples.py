from __future__ import annotations

from io import BytesIO
from pathlib import Path

import pytest
from pptx import Presentation

from packages.pipeline import generate
from packages.planner import GenerationInput

_EXAMPLES = sorted((Path(__file__).resolve().parent.parent / "examples").glob("*.json"))


@pytest.mark.parametrize("path", _EXAMPLES, ids=lambda p: p.stem)
def test_example_generates_valid_deck(path: Path) -> None:
    inp = GenerationInput.model_validate_json(path.read_text(encoding="utf-8"))
    res = generate(inp, with_preview=False)
    assert 25 <= len(res.deck_spec.slides) <= 30
    assert res.pptx_bytes is not None
    prs = Presentation(BytesIO(res.pptx_bytes))
    assert len(prs.slides) == len(res.deck_spec.slides)
