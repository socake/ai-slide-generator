from __future__ import annotations

from io import BytesIO

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from packages.asset_engine.engine import DEFAULT_INDEX_PATH
from packages.pipeline import generate
from packages.planner import GenerationInput


def test_placeholder_backgrounds_exist_and_open() -> None:
    base = DEFAULT_INDEX_PATH.parent / "backgrounds"
    for name in ("clean_blue.jpg", "exec_dark.jpg", "warm_gradient.jpg", "neutral_soft.jpg"):
        path = base / name
        assert path.exists()
        assert path.stat().st_size < 60_000  # 占位图必须很小
        Image.open(path).verify()


def test_pipeline_cover_gets_background_picture() -> None:
    res = generate(
        GenerationInput(topic="Python 入门", brief="变量、循环", audience="零基础"),
        with_preview=False,
    )
    assert res.pptx_bytes is not None
    cover = Presentation(BytesIO(res.pptx_bytes)).slides[0]
    assert any(sh.shape_type == MSO_SHAPE_TYPE.PICTURE for sh in cover.shapes)
